from pptx import Presentation
import sys

def analyze_ppt(file_path):
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i} ---")
        for j, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            print(f"  Shape {j}: {shape.text[:200]}")

if __name__ == '__main__':
    analyze_ppt(sys.argv[1])
